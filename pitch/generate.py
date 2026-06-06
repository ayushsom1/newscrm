"""Generate a sales-pitch .pptx for News CRM.

Re-runnable. The whole deck is constructed from primitives (rectangles,
textboxes) on blank layouts so we get pixel-level control without
fighting PowerPoint master themes.

    pip install python-pptx==1.0.2
    python pitch/generate.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
INK         = RGBColor(0x1C, 0x1A, 0x16)
BRAND       = RGBColor(0x9E, 0x1B, 0x17)
AI_BLUE     = RGBColor(0x25, 0x63, 0xEB)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE    = RGBColor(0xF7, 0xF6, 0xF4)
BORDER      = RGBColor(0xE5, 0xE3, 0xDE)
MUTED       = RGBColor(0x70, 0x6D, 0x66)
GREEN       = RGBColor(0x15, 0x80, 0x3D)
AMBER       = RGBColor(0xB4, 0x53, 0x09)
LIGHT_GREEN = RGBColor(0xE8, 0xF1, 0xEB)
LIGHT_AMBER = RGBColor(0xFD, 0xF2, 0xDF)
LIGHT_RED   = RGBColor(0xF9, 0xE3, 0xE3)
LIGHT_BLUE  = RGBColor(0xE6, 0xEE, 0xFD)

FONT = "Calibri"
FONT_HEAD = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, *, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    return s


def text(
    slide, left, top, width, height, content,
    *, font=FONT, size=14, bold=False, color=INK,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(
    slide, left, top, width, height, items,
    *, font=FONT, size=14, color=INK, marker="—",
    marker_color=None, spacing=8, line_spacing=1.2,
):
    if marker_color is None:
        marker_color = MUTED
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else spacing)
        p.line_spacing = line_spacing
        if marker:
            r1 = p.add_run()
            r1.text = marker + "  "
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.color.rgb = marker_color
            r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def chip(slide, left, top, label, *, color=INK, bg=OFFWHITE, border=BORDER, w=None):
    w = w or Inches(1.4)
    h = Inches(0.32)
    rect(slide, left, top, w, h, fill=bg, line=border)
    text(
        slide, left, top, w, h, label,
        font=FONT, size=10, bold=True, color=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def slide_footer(slide, n, total):
    text(
        slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.3),
        "News CRM · Sales overview",
        size=9, color=MUTED,
    )
    text(
        slide, Inches(12.0), Inches(7.0), Inches(0.9), Inches(0.3),
        f"{n} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL = 12


def add():
    return prs.slides.add_slide(BLANK)


# ===== Slide 1 — Cover =====================================================
s = add()
set_bg(s, INK)

# Brand red accent bar (top)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.18), fill=BRAND)

# Product mark
text(
    s, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
    "News CRM", size=14, bold=True, color=WHITE,
)

# Headline
text(
    s, Inches(0.7), Inches(2.4), Inches(11), Inches(1.4),
    "The AI-assisted CRM built for", size=44, bold=True, color=WHITE,
)
text(
    s, Inches(0.7), Inches(3.3), Inches(11), Inches(1.4),
    "regional newspapers.", size=44, bold=True, color=BRAND,
)

# Subline
text(
    s, Inches(0.7), Inches(4.7), Inches(11), Inches(1),
    "Two revenue engines. One operating principle: human on the loop.",
    size=18, color=RGBColor(0xCF, 0xCB, 0xC2),
)

# Footer
text(
    s, Inches(0.7), Inches(6.7), Inches(11), Inches(0.4),
    "Made for India and Nepal. Cloud or self-hosted.",
    size=11, color=RGBColor(0x9A, 0x95, 0x88), italic=True,
)


# ===== Slide 2 — Problem ===================================================
s = add()
set_bg(s, WHITE)

# Eyebrow + title
text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "THE PROBLEM", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "Regional dailies lose what they can't see.",
     size=32, bold=True, color=INK)

# Three problem cards
card_w = Inches(3.95)
card_h = Inches(3.6)
gap = Inches(0.2)
y = Inches(2.5)

problems = [
    (
        "Ad revenue softens silently",
        "Spend drops by 20–30% over two quarters before anyone catches it. By "
        "the time a renewal call goes out, the advertiser has already moved "
        "budget to digital.",
    ),
    (
        "Reminders and reviews are manual",
        "Renewal call lists are kept in spreadsheets. Complaint email goes "
        "to whoever is at the front desk. Billing disputes get handled the "
        "wrong way because nobody routes them.",
    ),
    (
        "Print runs over- or under-print",
        "Without a per-area returns model, circulation rounds up to be safe. "
        "5–15% of copies are wasted; the rest don't reach demand pockets.",
    ),
]

for i, (h, body) in enumerate(problems):
    x = Inches(0.6) + i * (card_w + gap)
    rect(s, x, y, card_w, card_h, fill=OFFWHITE, line=BORDER)
    # red corner accent
    rect(s, x, y, Inches(0.35), Inches(0.35), fill=BRAND)
    text(s, x + Inches(0.4), y + Inches(0.45), card_w - Inches(0.5),
         Inches(0.8), h, size=18, bold=True, color=INK)
    text(s, x + Inches(0.4), y + Inches(1.4), card_w - Inches(0.5),
         Inches(2.0), body, size=12, color=MUTED)

text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
     "All of this is solvable using data the newsroom already collects.",
     size=12, italic=True, color=INK)
slide_footer(s, 2, TOTAL)


# ===== Slide 3 — What we built =============================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "THE PRODUCT", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "One CRM. Two revenue engines.", size=32, bold=True, color=INK)
text(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5),
     "Sales, circulation and operations on the same screen — with AI that does "
     "the routine and a person who approves what matters.",
     size=14, color=MUTED)

pillars = [
    ("ADVERTISING", BRAND, [
        "Advertiser & contract CRUD",
        "Churn engine: spend trend + open rate + days to expiry",
        "AI proposal drafting with approval gate",
        "Government tender tracker",
    ]),
    ("CIRCULATION", AI_BLUE, [
        "Subscriber & subscription CRUD",
        "Renewal-risk engine + at-risk flagging",
        "Per-area print-run forecast",
        "Scheduled reminders on 14 / 7 / 3 / 1-day windows",
    ]),
    ("OPERATIONS", GREEN, [
        "AI complaint triage with engine fallback",
        "Auto-resolve routine ops; escalate sensitive ones",
        "Exception queue: a single pane for human attention",
        "Audit log of every AI and user action",
    ]),
]

pw = Inches(3.95)
ph = Inches(4.0)
py = Inches(2.7)
for i, (label, color, items) in enumerate(pillars):
    x = Inches(0.6) + i * (pw + Inches(0.2))
    rect(s, x, py, pw, ph, fill=WHITE, line=BORDER, line_w=Emu(9525))
    rect(s, x, py, pw, Inches(0.5), fill=color)
    text(s, x, py, pw, Inches(0.5), label,
         size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x + Inches(0.35), py + Inches(0.85), pw - Inches(0.5),
            ph - Inches(1.0), items, size=12, spacing=10)

slide_footer(s, 3, TOTAL)


# ===== Slide 4 — Human on the loop =========================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "THE DIFFERENCE", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "AI does the routine. A person approves anything that matters.",
     size=28, bold=True, color=INK)

zones = [
    (
        "AI HANDLES", GREEN, LIGHT_GREEN,
        "Routine work the AI can finish without you.",
        ["Non-delivery, pause, plan-change complaints",
         "Proposal first drafts grounded on real numbers",
         "Daily reminders for upcoming renewals"],
    ),
    (
        "HUMAN APPROVES", AMBER, LIGHT_AMBER,
        "AI prepares the work; you click Approve.",
        ["Every AI-drafted proposal before it is sent",
         "Each assistant 'proposed action' from the chat",
         "Anything flagged needs_human (e.g. high churn)"],
    ),
    (
        "ALWAYS ESCALATES", BRAND, LIGHT_RED,
        "Sensitive cases never touch the model.",
        ["Billing disputes, refunds, double charges",
         "Fraud, legal action, harassment, abuse",
         "Hard guardrail in the rules engine"],
    ),
]

zw = Inches(3.95)
zh = Inches(4.6)
zy = Inches(2.3)
for i, (label, accent, bg, sub, items) in enumerate(zones):
    x = Inches(0.6) + i * (zw + Inches(0.2))
    rect(s, x, zy, zw, zh, fill=bg, line=BORDER)
    text(s, x + Inches(0.4), zy + Inches(0.35), zw - Inches(0.7),
         Inches(0.5), label, size=12, bold=True, color=accent)
    text(s, x + Inches(0.4), zy + Inches(0.85), zw - Inches(0.7),
         Inches(0.7), sub, size=13, color=INK, bold=True)
    bullets(s, x + Inches(0.4), zy + Inches(1.85), zw - Inches(0.7),
            zh - Inches(2.2), items, size=12, spacing=8,
            marker="·", marker_color=accent)

text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
     "Every AI draft, triage and proposed action writes an AuditLog row.",
     size=11, italic=True, color=MUTED)


# ===== Slide 5 — Advertising side ==========================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "FOR SALES", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "Stop losing accounts to silent churn.",
     size=32, bold=True, color=INK)

# Left: deterministic engines
lx = Inches(0.6); ly = Inches(2.2)
lw = Inches(5.9); lh = Inches(4.6)
rect(s, lx, ly, lw, lh, fill=OFFWHITE, line=BORDER)
chip(s, lx + Inches(0.4), ly + Inches(0.4), "ENGINES", color=WHITE, bg=INK, border=INK)
text(s, lx + Inches(0.4), ly + Inches(0.95), lw - Inches(0.8), Inches(0.6),
     "Deterministic. No LLM. Fully tested.",
     size=18, bold=True, color=INK)
bullets(s, lx + Inches(0.4), ly + Inches(1.85), lw - Inches(0.8),
        lh - Inches(2.1),
        [
            "Churn engine: combines spend trend, proposal open rate and "
            "days-to-expiry into a 0-100 score and a low / med / high band.",
            "Pricing engine: Decimal money throughout, category multipliers, "
            "multi-day discounts, locale-aware GST or VAT.",
            "Re-scored nightly. Snapshots stored, never trusted.",
        ],
        size=13, spacing=10)

# Right: AI with approval
rx = Inches(6.8); ry = ly
rw = Inches(5.9); rh = lh
rect(s, rx, ry, rw, rh, fill=LIGHT_BLUE, line=BORDER)
chip(s, rx + Inches(0.4), ry + Inches(0.4), "AI + APPROVAL", color=WHITE, bg=AI_BLUE, border=AI_BLUE)
text(s, rx + Inches(0.4), ry + Inches(0.95), rw - Inches(0.8), Inches(0.6),
     "Drafts proposals, never sends them.",
     size=18, bold=True, color=INK)
bullets(s, rx + Inches(0.4), ry + Inches(1.85), rw - Inches(0.8),
        rh - Inches(2.1),
        [
            "AI drafter sees only the advertiser snapshot — name, contact, "
            "money as text, churn band. Prompt forbids inventing numbers.",
            "Every draft saved as DRAFT awaiting approval; high-churn "
            "accounts flagged needs_human and cannot be one-click approved.",
            "Send button is logged with actor, time and recipient.",
        ],
        size=13, spacing=10)

slide_footer(s, 5, TOTAL)


# ===== Slide 6 — Circulation side ==========================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "FOR CIRCULATION", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "Hold subscribers. Cut print waste.",
     size=32, bold=True, color=INK)

metrics = [
    ("At-risk flagging",
     "Renewal engine combines days-to-renew and missed payments into "
     "low / med / high severity, with explainable reasons on every row."),
    ("Print-run forecast",
     "Per-area target copies = (active + buffer) × (1 + returns %). "
     "Cuts the guesswork out of the print order."),
    ("Reminders that don't double-fire",
     "Daily job sends emails to advertisers and subscribers on 14 / 7 / 3 / 1-"
     "day windows. Idempotent per day; restart-safe."),
    ("Complaint triage built in",
     "Routine ops auto-resolve with action logged. Billing, fraud and "
     "harassment ALWAYS go to a human — guardrail fires before AI."),
]

mw = Inches(6.0); mh = Inches(1.95)
for i, (title, body) in enumerate(metrics):
    col, row = i % 2, i // 2
    x = Inches(0.6) + col * (mw + Inches(0.2))
    y = Inches(2.3) + row * (mh + Inches(0.2))
    rect(s, x, y, mw, mh, fill=WHITE, line=BORDER)
    rect(s, x, y, Inches(0.05), mh, fill=AI_BLUE)  # left accent bar
    text(s, x + Inches(0.35), y + Inches(0.3), mw - Inches(0.6),
         Inches(0.5), title, size=16, bold=True, color=INK)
    text(s, x + Inches(0.35), y + Inches(0.85), mw - Inches(0.6),
         mh - Inches(1.0), body, size=12, color=MUTED)

slide_footer(s, 6, TOTAL)


# ===== Slide 7 — Exception queue ===========================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "THE DASHBOARD", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "One screen for everything that needs a human.",
     size=28, bold=True, color=INK)
text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
     "The exception queue is derived live from the CRM — it cannot drift from "
     "reality. Each row tells you exactly where to click.",
     size=13, color=MUTED)

# Mock queue
qx = Inches(0.6); qy = Inches(2.7)
qw = Inches(12.1); qh = Inches(4.0)
rect(s, qx, qy, qw, qh, fill=WHITE, line=BORDER)
# header band
rect(s, qx, qy, qw, Inches(0.45), fill=OFFWHITE, line=BORDER)
text(s, qx + Inches(0.3), qy, Inches(6), Inches(0.45),
     "Exception queue",
     size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# severity tab pills
tab_y = qy + Inches(0.08); tab_x = qx + Inches(8.7)
for label, color in (("HUMAN  3", BRAND), ("APPROVE  4", AMBER), ("AI HANDLED  12", GREEN)):
    rect(s, tab_x, tab_y, Inches(1.05), Inches(0.3), fill=WHITE, line=color)
    text(s, tab_x, tab_y, Inches(1.05), Inches(0.3), label,
         size=9, bold=True, color=color, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    tab_x += Inches(1.1)

rows = [
    ("HUMAN",   BRAND,  LIGHT_RED,
     "Sunrise Builders flagged high churn",   "score=76  ·  spend down -55%"),
    ("HUMAN",   BRAND,  LIGHT_RED,
     "Harish Chand Verma — 3 missed payments", "area=Saket Nagar"),
    ("APPROVE", AMBER,  LIGHT_AMBER,
     "Renewal proposal awaiting approval — Greenleaf Restaurant",
     "source=AI_DRAFT  ·  contract expires in 7d"),
    ("APPROVE", AMBER,  LIGHT_AMBER,
     "Public Health — Dengue Awareness Campaign",
     "deadline in 7d  ·  est ₹2,50,000"),
    ("AI HANDLED", GREEN, LIGHT_GREEN,
     "Wet copy complaint — Rashmi Agarwal",
     "auto-resolved  ·  replacement scheduled"),
    ("AI HANDLED", GREEN, LIGHT_GREEN,
     "Pause request — Meera Iyer",
     "auto-resolved  ·  paused 15-30 of month"),
]

row_h = Inches(0.5)
row_y = qy + Inches(0.5)
for label, accent, bg, summary, detail in rows:
    rect(s, qx, row_y, qw, row_h, fill=bg, line=BORDER)
    # severity badge
    chip(s, qx + Inches(0.2), row_y + Inches(0.1), label,
         color=WHITE, bg=accent, border=accent, w=Inches(1.0))
    text(s, qx + Inches(1.4), row_y, Inches(8.0), row_h, summary,
         size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, qx + Inches(9.4), row_y, Inches(2.6), row_h, detail,
         size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    row_y += row_h + Emu(0)

slide_footer(s, 7, TOTAL)


# ===== Slide 8 — Architecture =============================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "HOW IT'S BUILT", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "Split the brain. So the AI can never get it wrong.",
     size=28, bold=True, color=INK)

# Engines column
ex = Inches(0.6); ey = Inches(2.3); ew = Inches(6.0); eh = Inches(3.4)
rect(s, ex, ey, ew, eh, fill=OFFWHITE, line=BORDER)
text(s, ex + Inches(0.4), ey + Inches(0.4), ew, Inches(0.4),
     "ENGINES", size=11, bold=True, color=INK)
text(s, ex + Inches(0.4), ey + Inches(0.85), ew - Inches(0.8), Inches(0.6),
     "Arithmetic. Rules. Money.", size=20, bold=True, color=INK)
bullets(s, ex + Inches(0.4), ey + Inches(1.7), ew - Inches(0.8),
        eh - Inches(2.0),
        ["Pricing · churn · renewal · print-run · triage",
         "Pure functions. No LLM. Decimal end-to-end.",
         "Unit-tested boundary by boundary. Repeatable."],
        size=13, spacing=8)

# AI column
ax = Inches(6.8); aw = Inches(6.0)
rect(s, ax, ey, aw, eh, fill=LIGHT_BLUE, line=BORDER)
text(s, ax + Inches(0.4), ey + Inches(0.4), aw, Inches(0.4),
     "AI (OpenRouter)", size=11, bold=True, color=AI_BLUE)
text(s, ax + Inches(0.4), ey + Inches(0.85), aw - Inches(0.8), Inches(0.6),
     "Language only.", size=20, bold=True, color=INK)
bullets(s, ax + Inches(0.4), ey + Inches(1.7), aw - Inches(0.8),
        eh - Inches(2.0),
        ["Drafts proposals · triages complaints · grounded chat",
         "Pydantic-validated output. Engine fallback on any failure.",
         "Swap model via one env var: GPT, Claude, Llama, Gemini."],
        size=13, spacing=8)

# Audit band
bx = Inches(0.6); by = Inches(5.95); bw = Inches(12.1); bh = Inches(0.95)
rect(s, bx, by, bw, bh, fill=INK)
text(s, bx + Inches(0.4), by, bw - Inches(0.8), bh,
     "AUDIT LOG  ·  every AI draft, every triage, every approve / reject / send writes a row",
     size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

slide_footer(s, 8, TOTAL)


# ===== Slide 9 — Autonomy dial =============================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "CONTROL", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "Operate the AI on your terms.",
     size=32, bold=True, color=INK)

# Settings mock
sx = Inches(0.6); sy = Inches(2.3); sw = Inches(7.5); sh = Inches(4.5)
rect(s, sx, sy, sw, sh, fill=WHITE, line=BORDER)
rect(s, sx, sy, sw, Inches(0.45), fill=OFFWHITE)
text(s, sx + Inches(0.3), sy, Inches(6), Inches(0.45),
     "Settings  ·  Autonomy", size=11, bold=True, color=INK,
     anchor=MSO_ANCHOR.MIDDLE)

toggles = [
    ("AI complaint triage",                   True),
    ("Auto-resolve routine triage",           True),
    ("AI proposal drafting",                  True),
    ("High-churn drafts always need a human", True),
    ("Assistant actions require ADMIN",       False),
]
ty = sy + Inches(0.7)
for label, on in toggles:
    rect(s, sx + Inches(0.25), ty, sw - Inches(0.5), Inches(0.6),
         fill=WHITE, line=BORDER)
    text(s, sx + Inches(0.5), ty, Inches(5.5), Inches(0.6), label,
         size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    # toggle pill
    pill_x = sx + sw - Inches(1.05)
    pill_bg = INK if on else RGBColor(0xCF, 0xCB, 0xC2)
    rect(s, pill_x, ty + Inches(0.18), Inches(0.7), Inches(0.24),
         fill=pill_bg, line=pill_bg)
    knob_x = pill_x + (Inches(0.45) if on else Inches(0.02))
    rect(s, knob_x, ty + Inches(0.21), Inches(0.18), Inches(0.18),
         fill=WHITE, line=WHITE)
    ty += Inches(0.72)

# Right column — narrative
nx = Inches(8.4); ny = sy; nw = Inches(4.3); nh = sh
rect(s, nx, ny, nw, nh, fill=OFFWHITE, line=BORDER)
text(s, nx + Inches(0.3), ny + Inches(0.4), nw, Inches(0.4),
     "WHAT THE DIAL DOES", size=10, bold=True, color=BRAND)
bullets(s, nx + Inches(0.3), ny + Inches(0.95), nw - Inches(0.5),
        nh - Inches(1.1),
        ["Take AI off a feature and the engine fallback runs immediately.",
         "Force everything to a human while you build trust in a pilot.",
         "Switch models via env: openai/gpt-4o-mini, anthropic/*, llama, gemini.",
         "Locale toggle changes currency and tax label across the whole CRM.",
         "Audit log records who flipped which dial and when."],
        size=12, spacing=10)

slide_footer(s, 9, TOTAL)


# ===== Slide 10 — Locale ===================================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "BUILT FOR YOUR MARKET", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "India and Nepal, day one.",
     size=32, bold=True, color=INK)

# India
ix = Inches(0.6); iy = Inches(2.3); iw = Inches(6.0); ih = Inches(4.2)
rect(s, ix, iy, iw, ih, fill=OFFWHITE, line=BORDER)
rect(s, ix, iy, iw, Inches(0.6), fill=BRAND)
text(s, ix, iy, iw, Inches(0.6), "INDIA",
     size=14, bold=True, color=WHITE,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, ix + Inches(0.45), iy + Inches(1.0), iw - Inches(0.8),
        ih - Inches(1.2),
        ["Currency ₹ INR  ·  5% GST on classifieds",
         "Sunday weekly special supported",
         "Hindi (Devanagari) and English in UI and AI output",
         "DIPR / government tender tracker out of the box"],
        size=14, spacing=12)

# Nepal
nx = Inches(6.85); ny = iy; nw = iw; nh = ih
rect(s, nx, ny, nw, nh, fill=OFFWHITE, line=BORDER)
rect(s, nx, ny, nw, Inches(0.6), fill=AI_BLUE)
text(s, nx, ny, nw, Inches(0.6), "NEPAL",
     size=14, bold=True, color=WHITE,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, nx + Inches(0.45), ny + Inches(1.0), nw - Inches(0.8),
        nh - Inches(1.2),
        ["Currency NPR  ·  13% VAT on classifieds",
         "Saturday weekly special supported",
         "Nepali (Devanagari) and English in UI and AI output",
         "Pricing tiers and tax rules configured per locale"],
        size=14, spacing=12)

text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
     "Commercial rules live in Settings, not in code. Switch locale, see the "
     "whole CRM re-render.",
     size=11, italic=True, color=MUTED)

slide_footer(s, 10, TOTAL)


# ===== Slide 11 — Pilot path ===============================================
s = add()
set_bg(s, WHITE)

text(s, Inches(0.6), Inches(0.55), Inches(6), Inches(0.3),
     "PILOT PATH", size=10, bold=True, color=BRAND)
text(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
     "An 8-week pilot, then a clear decision.",
     size=32, bold=True, color=INK)

phases = [
    ("WEEKS 1-2", "Onboard",
     "Import advertisers and subscribers. Train sales and circulation staff. "
     "AI off; engines computing."),
    ("WEEKS 3-4", "Engines online",
     "Churn, renewal and forecast live. AI in suggest-only mode for drafts; "
     "nothing auto-sent."),
    ("WEEKS 5-6", "AI assist",
     "Triage and drafting on with engine fallback. Reminders firing. "
     "Audit log under review."),
    ("WEEKS 7-8", "Decide",
     "Dial autonomy to the level you trust. Review revenue impact. "
     "Continue, pause or expand."),
]

bw = Inches(2.95); bh = Inches(3.6)
by = Inches(2.4)
for i, (period, name, body) in enumerate(phases):
    x = Inches(0.6) + i * (bw + Inches(0.2))
    rect(s, x, by, bw, bh, fill=WHITE, line=BORDER)
    rect(s, x, by, bw, Inches(0.45), fill=INK)
    text(s, x, by, bw, Inches(0.45), period,
         size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.3), by + Inches(0.7), bw - Inches(0.6),
         Inches(0.5), name, size=18, bold=True, color=INK)
    text(s, x + Inches(0.3), by + Inches(1.4), bw - Inches(0.6),
         bh - Inches(1.6), body, size=12, color=MUTED)

text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
     "All projected metrics are illustrative targets to validate in the pilot — "
     "never promised on day one.",
     size=11, italic=True, color=MUTED)

slide_footer(s, 11, TOTAL)


# ===== Slide 12 — CTA ======================================================
s = add()
set_bg(s, INK)

rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.18), fill=BRAND)

text(s, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
     "News CRM", size=14, bold=True, color=WHITE)

text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2),
     "Let's run a pilot on your top 50 advertisers.",
     size=36, bold=True, color=WHITE)
text(s, Inches(0.7), Inches(3.4), Inches(12), Inches(0.6),
     "Two weeks. Your data. No commitment.",
     size=22, color=RGBColor(0xCF, 0xCB, 0xC2))

# CTA box
cx = Inches(0.7); cy = Inches(4.6); cw = Inches(11.9); ch = Inches(1.7)
rect(s, cx, cy, cw, ch, fill=RGBColor(0x2A, 0x27, 0x21), line=BRAND)
text(s, cx + Inches(0.5), cy + Inches(0.25), cw, Inches(0.4),
     "NEXT STEPS", size=10, bold=True, color=BRAND)
bullets(s, cx + Inches(0.5), cy + Inches(0.75), cw - Inches(1.0),
        Inches(0.9),
        ["30-minute product walkthrough with your sales and circulation leads",
         "Two-week sandbox with your top 50 advertisers and one delivery area",
         "Joint review and a go / no-go at the end"],
        size=13, color=WHITE, marker="·", marker_color=BRAND, spacing=4)

text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
     "contact@news-crm.example     ·     news-crm.example",
     size=12, color=RGBColor(0xCF, 0xCB, 0xC2))


# ---------------------------------------------------------------------------
# Write
# ---------------------------------------------------------------------------
OUT = Path(__file__).parent / "news-crm-pitch.pptx"
prs.save(OUT)
print(f"deck written: {OUT}")
